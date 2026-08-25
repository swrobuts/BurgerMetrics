#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
bau_datenmodell.py — Foliendeck "BurgerMetrics: Datenmodell und Aufbau".

Dokumentiert, wie die drei Anwendungen des Projekts zusammenhaengen und warum
sie drei verschiedene Datenmodelle brauchen: Shop, Warenwirtschaft, Auswertung.
Das Schlusskapitel ordnet den gewaehlten Weg in die Alternativen ein — Data
Warehouse, Data Lake, Lakehouse, Data Mesh.

Die Diagramme stammen aus diagramme/*.mmd und werden von render_mermaid.mjs
erzeugt — dieses Skript setzt sie nur.

    python3 bau_datenmodell.py                  # -> ../../BurgerMetrics_Datenmodell.pptx
    python3 bau_datenmodell.py -o /pfad.pptx

Voraussetzungen: pip install python-pptx, Skill thws-slides (THWS_SKILL),
gerenderte PNG in diagramme/ (node render_mermaid.mjs).

Alle Messwerte im Deck stammen aus einem Lauf gegen den Datenbestand in
dataset/ (DuckDB 1.5.5, Median aus sieben Wiederholungen nach Warmlauf).
"""
import argparse
import importlib.util
import os
import struct
import sys
import types
from pathlib import Path


def finde_skill():
    if os.environ.get("THWS_SKILL"):
        p = Path(os.environ["THWS_SKILL"])
        if (p / "assets" / "template.pptx").exists():
            return p
        sys.exit(f"FEHLER: THWS_SKILL zeigt auf {p}, dort fehlt assets/template.pptx")
    basis = Path.home() / "Library/Application Support/Claude"
    treffer = sorted(basis.glob("**/skills/thws-slides/assets/template.pptx"))
    if treffer:
        return treffer[-1].parent.parent
    sys.exit('FEHLER: Skill thws-slides nicht gefunden. export THWS_SKILL="/Pfad/zu/skills/thws-slides"')


SKILL = finde_skill()
SPECS = SKILL / "assets" / "specs"


def _load_spec(name="bint"):
    sp = importlib.util.spec_from_file_location(f"spec_{name}", SPECS / f"{name}.py")
    m = importlib.util.module_from_spec(sp)
    sp.loader.exec_module(m)
    return m


# spec_loader des Skills zeigt auf <skill>/specs statt <skill>/assets/specs.
_fake = types.ModuleType("spec_loader")
_fake.load = _load_spec
_fake.SPECS_DIR = SPECS
sys.modules["spec_loader"] = _fake

sys.path.insert(0, str(SKILL / "scripts"))
import bausteine as B                                    # noqa: E402
from pptx import Presentation                            # noqa: E402
from pptx.util import Pt                                 # noqa: E402
from pptx.enum.text import PP_ALIGN                      # noqa: E402

G = B.G
BLUE, WARN, GOOD, METHOD, HINT = G.BLUE, G.WARN, G.GOOD, G.METHOD, G.HINT
CL, CW, CR = G.CONTENT_LEFT, G.CONTENT_WIDTH, G.CONTENT_RIGHT
Y0, YMAX = G.CONTENT_ZONE_Y_MIN, G.CONTENT_ZONE_Y_MAX
HOEHE = YMAX - Y0


# --- Absatzformat: as_paras() erkennt nur Listen von LISTEN als mehrere Absaetze.
#     Flache Stringlisten liefen sonst zu einem Absatz zusammen.
def _paras(x):
    if isinstance(x, list) and x and all(isinstance(e, str) for e in x):
        return [[e] for e in x]
    return x


_kachel, _band, _gegenueber = B.kachel, B.band, B.gegenueber
B.kachel = lambda s, x, y, w, h, f, t, b, gap=5: _kachel(s, x, y, w, h, f, t, _paras(b), gap)
B.band = lambda s, y, h, p, **k: _band(s, y, h, _paras(p), **k)
B.gegenueber = lambda s, y, h, l, r, **k: _gegenueber(
    s, y, h, (l[0], l[1], _paras(l[2])), (r[0], r[1], _paras(r[2])), **k)

ap = argparse.ArgumentParser(description="Foliendeck Datenmodell bauen")
ap.add_argument("-o", "--ausgabe", default=None)
args = ap.parse_args()

HIER = Path(__file__).resolve().parent
DIA = HIER / "diagramme"
TEMPLATE = SKILL / "assets" / "template.pptx"
OUT = Path(args.ausgabe) if args.ausgabe else HIER.parent.parent / "BurgerMetrics_Datenmodell.pptx"

QUELLE = "Eigene Darstellung · Diagramme aus slides/diagramme/*.mmd"
Q_MESS = ("Eigene Messung · DuckDB 1.5.5 auf dataset/, Median aus sieben Läufen · "
          "Diagramme aus slides/diagramme/*.mmd")
Q_LIT = ("Eigene Darstellung nach Armbrust u. a., Lakehouse, CIDR 2021; "
         "Marktzahlen 2026 aus Anbieter- und Beraterberichten")
Q_MESH = ("Eigene Darstellung nach Dehghani, Data Mesh, 2022; "
          "Reifegrad- und Erfolgszahlen aus Beraterberichten 2025/2026")
Q_EIG = "Eigene Einschätzung · Diagramme aus slides/diagramme/*.mmd"
Q_DUCK = "Eigene Darstellung nach dem DuckLake-Manifest, ducklake.select, Stand 2026"

prs = Presentation(str(TEMPLATE))
LAYOUT = {l.name: l for l in prs.slide_layouts}


def neu(name):
    return prs.slides.add_slide(LAYOUT[name])


def png_groesse(p):
    d = open(p, "rb").read(33)
    return struct.unpack(">II", d[16:24])


def bild(slide, name, y=None, max_h=None, x=None, max_w=None):
    """Diagramm linksbuendig auf der Inhaltsflucht, seitenverhaeltnistreu.

    Der Schriftgrad im Bild folgt aus dem Massstab: Mermaid setzt 15 CSS-px,
    die PNG sind dreifach ueberabgetastet. Auf der Folie kommen davon
    45 * Massstab Punkte an — unter 11 pt ist das im Hoersaal nicht lesbar.
    """
    p = DIA / f"{name}.png"
    if not p.exists():
        sys.exit(f"FEHLER: {p.name} fehlt — erst 'node render_mermaid.mjs' ausfuehren.")
    pw, ph = png_groesse(p)
    y = Y0 if y is None else y
    max_h = (YMAX - y) if max_h is None else max_h
    max_w = CW if max_w is None else max_w
    sk = min(max_w / pw, max_h / ph)
    if 45 * sk < 10.5:
        B.WARNINGS.append(f"Diagramm {name}: Beschriftung nur {45 * sk:.1f} pt")
    w, h = pw * sk, ph * sk
    slide.shapes.add_picture(str(p), B.E(CL if x is None else x), B.E(y), B.E(w), B.E(h))
    return w, h


def einl(slide, text, variante="Lisa_Slide"):
    ph = B.einleitung(slide, [text], variante)
    vor = [q for q in LAYOUT[variante].placeholders if q.placeholder_format.idx == 14][0]
    ph.left, ph.top, ph.width = vor.left, vor.top, vor.width
    ph.height = Pt(63.8)
    for p in ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.JUSTIFY
        for r in p.runs:
            r.font.size = Pt(14)
    return ph


def aufraeumen(slide):
    for sh in list(slide.placeholders):
        try:
            if not sh.text_frame.text.strip():
                sh._element.getparent().remove(sh._element)
        except Exception:
            pass


def kapitel(titel):
    s = neu("Chapter")
    for sh in s.placeholders:
        if str(sh.placeholder_format.type).startswith("TITLE"):
            sh.text_frame.text = titel
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(29)
    aufraeumen(s)


def bewertung(slide, y, spalten, zeilen, rh=44.0, bw=250.0):
    """Bewertungsmatrix: Zeilenbezeichner links, Balken je Spalte.

    Balken statt Harvey-Kreisen: Ein teilgefuellter Kreis ist eine PIE-Form,
    deren Winkel LibreOffice beim Rendern verwirft — halbvolle Kreise kamen
    dort als leere oder volle heraus. Rechtecke stellt jedes Programm gleich dar.

    spalten: Liste von Spaltenueberschriften.
    zeilen:  Liste von (Bezeichner, [Anteil 0..1 je Spalte], Hervorhebung).
    """
    n = len(spalten)
    sw = (CR - CL - bw) / n
    balken_w = min(100.0, sw - 34)
    for j, kopfzeile in enumerate(spalten):
        B.textbox(slide, CL + bw + j * sw, y, sw, 30, [(kopfzeile, False, B.SEC)],
                  10.5, align=PP_ALIGN.CENTER)
    yy = y + 34
    for bez, werte, hl in zeilen:
        farbe = BLUE if hl else B.SEC
        B.shape(slide, B.MSO_SHAPE.RECTANGLE, CL, yy, bw + n * sw, rh,
                "F4F7FA" if hl else B.WHITE, B.HAIR)
        B.textbox(slide, CL + 14, yy + 6, bw - 20, rh - 12, [(bez, hl, farbe)],
                  12, anchor=B.MSO_ANCHOR.MIDDLE)
        for j, a in enumerate(werte):
            bx = CL + bw + j * sw + (sw - balken_w) / 2
            by = yy + rh / 2 - 4
            B.shape(slide, B.MSO_SHAPE.RECTANGLE, bx, by, balken_w, 8, "E8E6E0")
            if a > 0:
                B.shape(slide, B.MSO_SHAPE.RECTANGLE, bx, by, balken_w * a, 8, farbe)
        yy += rh + 6
    return yy


# ═══════════════════════════════════════════════════ Deckblatt
s = neu("Frontpage_Digital")
for sh in s.placeholders:
    i = sh.placeholder_format.idx
    if i == 10:
        sh.text_frame.text = "BurgerMetrics"
    elif i == 11:
        sh.text_frame.text = "Datenmodell und Aufbau"
aufraeumen(s)

# ═══════════════════════════════════════════════════ Der rote Faden
s = neu("Lisa_Slide")
B.kopf(s, "Leitfrage", "Ein Datenmodell ist eine Antwort — und sie hat einen Preis", QUELLE)
einl(s, "Dieses Deck folgt einer einzigen Frage: Welche Frage stellt das System, und "
        "welches Modell beantwortet sie am günstigsten? Günstig heißt nicht nur schnell. "
        "Jede Modellentscheidung verschiebt Kosten zwischen drei Konten, und keine "
        "Entscheidung senkt alle drei zugleich.")
B.chevron_kette(s, Y0, 38, ["1 · Drei Systeme, drei Fragen",
                            "2 · Konsistenz zuerst",
                            "3 · Lesetempo zuerst",
                            "4 · Ein Weg von vielen"])
w3 = (CW - 2 * 18) / 3
for i, (t, b, c) in enumerate([
    ("Schreibkosten", ["Was kostet es, einen Sachverhalt widerspruchsfrei zu halten?",
                       "Steigt mit jeder Stelle, an der dasselbe Merkmal steht.",
                       "Zahlt das operative System."], BLUE),
    ("Lesekosten", ["Was kostet eine Auswertung über den ganzen Bestand?",
                    "Steigt mit der Zahl der Verknüpfungen je Abfrage.",
                    "Zahlt das Auswertungssystem."], GOOD),
    ("Betriebskosten", ["Was kostet der Apparat, der beides am Laufen hält?",
                        "Steigt mit jeder zusätzlichen Komponente und Zuständigkeit.",
                        "Zahlt die Organisation — Thema von Kapitel 4."], METHOD),
]):
    B.kachel(s, CL + i * (w3 + 18), Y0 + 56, w3, 172, c, t, b)
B.band(s, Y0 + 242, 66, [
    "Das operative Modell nimmt hohe Lesekosten in Kauf, um Schreibkosten zu sparen. Das "
    "Auswertungsmodell dreht das um. Kapitel 4 zeigt Architekturen, die vor allem an der "
    "dritten Kostenart drehen — und was sie dafür verlangen."])

# ═══════════════════════════════════════════════════ Teil 1
kapitel("Ein Geschäftsvorfall, drei Systeme")

# Systemkontext
s = neu("Lisa_Slide")
B.kopf(s, "Überblick", "Eine Bestellung durchläuft drei Systeme mit drei Aufgaben", QUELLE)
einl(s, "Wer eine Bestellung aufgibt, löst eine Kette aus: Der Shop nimmt sie entgegen, die "
        "Warenwirtschaft verbucht sie, die Auswertung zählt sie später mit. Im Datenbestand "
        "sind das 754.513 Bestellungen aus neun Jahren, acht Filialen und zwei Kanälen — "
        "erfasst zwischen dem 15. März 2017 und dem 31. März 2026.")
bild(s, "01_systemkontext", y=Y0 + 10, max_h=175)
B.band(s, Y0 + 200, 62, [
    "Der gestrichelte Pfeil ist die entscheidende Stelle: Die Auswertung liest nicht aus dem "
    "operativen System, sondern aus einem periodisch erzeugten Abzug. Sonst würde eine "
    "Jahresauswertung über 3,7 Millionen Zeilen den laufenden Kassenbetrieb ausbremsen."])

# Shop-Modell
s = neu("Lisa_Slide")
B.kopf(s, "System 1 · Website", "Der Shop speichert nur, was bis zum Kaufabschluss gebraucht wird", QUELLE)
einl(s, "Ein Warenkorb ist kein Beleg. Er lebt in einer Sitzung, ändert sich mit jedem Klick "
        "und verschwindet, wenn nichts daraus wird. Entsprechend schmal ist das Modell: drei "
        "Tabellen für Sitzung, Positionen und Artikelstamm. Ein Kunde ist optional — "
        "bestellen kann man auch ohne Anmeldung.")
bw, _ = bild(s, "02_shop_modell", max_w=470)
B.kachel(s, CL + bw + 24, Y0, CR - (CL + bw + 24), 190, METHOD, "Worauf dieses Modell optimiert ist",
         ["Sehr viele kleine Schreibvorgänge: jeder Klick ändert den Warenkorb.",
          "Kurze Lebensdauer — abgebrochene Sitzungen werden verworfen.",
          "Keine Historie, kein Bezug zu früheren Käufen.",
          "Mit dem Kaufabschluss übergibt der Shop an die Warenwirtschaft und ist fertig."])

# ═══════════════════════════════════════════════════ Teil 2
kapitel("Das operative Modell: Konsistenz zuerst")

# Verkaufspfad
s = neu("Lisa_Slide")
B.kopf(s, "System 2 · Warenwirtschaft", "Aus dem Warenkorb wird ein Beleg mit festen Bezügen", QUELLE)
einl(s, "Mit dem Kaufabschluss wird aus einem flüchtigen Warenkorb ein Beleg, der Jahre "
        "aufbewahrt wird. Er braucht feste Bezüge: welche Filiale, welcher Kunde, welche "
        "Zahlungsart. Die Positionen hängen zwingend am Kopf — im Bestand sind es im Mittel "
        "3,91 Positionen je Bestellung, nie null.")
bw, _ = bild(s, "03_wawi_verkaufspfad", max_w=470, max_h=295)
B.kachel(s, CL + bw + 24, Y0, CR - (CL + bw + 24), 190, BLUE, "Die Notation lesen",
         ["Doppelstrich = genau eins, Krähenfuß = viele.",
          "kundenbestellung ||--|{ bestellposition: eine Bestellung hat mindestens eine Position.",
          "filiale ||--o{ kundenbestellung: eine Filiale kann auch null Bestellungen haben.",
          "Kreis = optional, Strich = verpflichtend."])

# Bereiche
s = neu("Lisa_Slide")
B.kopf(s, "System 2 · Warenwirtschaft", "Der Verkauf ist ein Bereich von sieben — 26 Tabellen insgesamt", QUELLE)
einl(s, "Ein Warenwirtschaftssystem bildet den ganzen Betrieb ab, nicht nur den Verkauf: "
        "26 Tabellen in sieben Bereichen, jedes Merkmal genau einmal. Für die Absatzfragen "
        "der Auswertung werden davon drei Bereiche gebraucht — Stammdaten, Verkauf und die "
        "externen Wetterdaten.")
bild(s, "04_wawi_bereiche", y=Y0, max_h=185)
wb = (CW - 18) / 2
B.kachel(s, CL, Y0 + 200, wb, 122, BLUE, "Warum so viele Tabellen",
         ["Einkauf, Lager und Personal hängen am selben Artikelstamm.",
          "Jedes Merkmal steht genau einmal — das ist die dritte Normalform."])
B.kachel(s, CL + wb + 18, Y0 + 200, wb, 122, HINT, "Was davon die Auswertung braucht",
         ["Einkauf, Lager, Personal und Rechnungswesen beantworten keine Absatzfrage.",
          "Aus 26 Tabellen werden dadurch zwölf."])

# Normalisierung
s = neu("Lisa_Slide")
B.kopf(s, "System 2 · Warenwirtschaft", "Drei Tabellen für einen Artikel — und warum das richtig ist", QUELLE)
einl(s, "Ein Artikel hat eine Unterkategorie, die zu einer Kategorie gehört. Im operativen "
        "Modell sind das drei Tabellen. Das wirkt umständlich, hat aber einen Grund: Jeder "
        "Kategoriename steht genau einmal. Wer ihn ändert, ändert eine Zeile — nicht "
        "siebenundfünfzig, so viele Produkte führt der Bestand.")
B.gegenueber(s, Y0, 190,
             (GOOD, "Was Normalisierung leistet",
              ["Kein Merkmal steht doppelt, also kann nichts widersprüchlich werden.",
               "Änderungen treffen genau eine Zeile.",
               "Einfüge- und Löschanomalien sind ausgeschlossen: Eine Kategorie kann angelegt werden, bevor es Artikel gibt."]),
             (HINT, "Was sie kostet",
              ["Jede Auswertung muss die Tabellen wieder zusammenführen.",
               "Umsatz je Kategorie braucht vier Tabellen und drei Verknüpfungen.",
               "Bei Millionen Zeilen und vielen gleichzeitigen Abfragen wird das spürbar."]),
             badge_l="+", badge_r="−")

# ═══════════════════════════════════════════════════ Teil 3
kapitel("Das Auswertungsmodell: Lesetempo zuerst")

# Denormalisierung
s = neu("Lisa_Slide")
B.kopf(s, "System 3 · Auswertung", "Für die Auswertung wird die Normalisierung zurückgenommen", QUELLE)
einl(s, "Das Auswertungsmodell dreht die Abwägung um. Die drei Artikel-Tabellen werden zu "
        "einer breiten Dimensionstabelle zusammengezogen. Kategorie und Unterkategorie stehen "
        "danach redundant in jeder der 57 Produktzeilen — und genau das ist gewollt.")
bild(s, "05_denormalisierung", y=Y0, max_h=130)
B.kachel(s, CL, Y0 + 148, CW, 112, METHOD, "Warum die Redundanz vertretbar ist",
         ["Ein Auswertungsbestand wird periodisch neu beladen, nicht laufend fortgeschrieben.",
          "Änderungsanomalien entstehen im Schreibbetrieb — den gibt es hier nicht.",
          "Aus drei Verknüpfungen je Abfrage wird eine."])
B.band(s, Y0 + 274, 48, [
    "Faustregel: Konsistenz hat Vorrang im operativen System, Lesegeschwindigkeit im Auswertungssystem."])

# Granularität
s = neu("Lisa_Slide")
B.kopf(s, "System 3 · Auswertung", "Die erste Frage an jede Faktentabelle: Welches Ereignis ist eine Zeile?", QUELLE)
einl(s, "Bestellung 1 aus dem Datenbestand, vollständig. Der Bestellkopf trägt den Rabatt und "
        "den Endbetrag, die drei Positionen tragen Menge und Einzelpreis. Beides sind Fakten — "
        "aber auf verschiedenen Ebenen. Diese Unterscheidung heißt Granularität.")
bw, _ = bild(s, "06_granularitaet", max_w=380)
B.sprechblase(s, CL + bw + 30, Y0, CR - (CL + bw + 30), 196, _paras([
    [("Warum das nicht in eine Tabelle passt", True, WARN)],
    ["Läge alles auf Positionsebene, stünde der Rabatt dreimal da — und jede Summe darüber wäre das Dreifache."],
    ["Läge alles auf Bestellebene, gäbe es keine Produktanalysen mehr."],
    ["Die Positionssummen ergeben hier 10,35 € und damit den Bruttobetrag, nicht den Nettobetrag von 8,80 €."],
]), farbe=WARN)

# Galaxy
s = neu("Lisa_Slide")
B.kopf(s, "System 3 · Auswertung", "Zwei Granularitäten verlangen zwei Faktentabellen", QUELLE)
einl(s, "Ein Stern-Schema hat definitionsgemäß eine Faktentabelle. Zwei Faktentabellen, die "
        "sich Dimensionen teilen, heißen Galaxy-Schema oder Fact Constellation — die "
        "naheliegende Erweiterung, sobald ein Sachverhalt auf zwei Ebenen gemessen wird. "
        "dim_product hängt nur an der Positionsebene, die übrigen Dimensionen an beiden.")
bild(s, "07_galaxy", y=Y0, max_h=HOEHE)

# Die drei Modelle
s = neu("Lisa_Slide")
B.kopf(s, "Zusammenschau", "Drei Modelle, weil drei verschiedene Fragen gestellt werden", QUELLE)
einl(s, "Es gibt nicht ein richtiges Datenmodell für BurgerMetrics, sondern drei — je eines "
        "für die Aufgabe, die das System zu erfüllen hat. Wer das Auswertungsmodell in die "
        "Kasse einbaut, bekommt Inkonsistenzen; wer die Kasse auswertet, bremst den Betrieb.")
for i, (t, b, c) in enumerate([
    ("Shop · Website", ["Frage: Was liegt gerade im Warenkorb?",
                        "viele kleine Schreibvorgänge",
                        "kurze Lebensdauer, keine Historie",
                        "3 Tabellen"], METHOD),
    ("Warenwirtschaft", ["Frage: Was ist verbindlich passiert?",
                         "Konsistenz vor Geschwindigkeit",
                         "dritte Normalform, keine Redundanz",
                         "26 Tabellen"], BLUE),
    ("Auswertung", ["Frage: Wie entwickelt sich das Geschäft?",
                    "Lesen über 3,7 Millionen Zeilen",
                    "bewusste Redundanz, wenige Verknüpfungen",
                    "12 Tabellen"], GOOD),
]):
    B.kachel(s, CL + i * (w3 + 18), Y0, w3, 168, c, t, b)

# Kette
s = neu("Lisa_Slide")
B.kopf(s, "Umsetzung", "Vom Quellsystem zum Bericht in drei Schichten", Q_MESS)
einl(s, "Der Umbau geschieht nicht in einem Sprung, sondern in Schichten. Die Rohdaten bleiben "
        "unangetastet, damit jederzeit nachvollziehbar ist, was das Quellsystem geliefert hat. "
        "Erst die mittlere Schicht typisiert und benennt, erst die dritte modelliert um.")
bild(s, "08_kette", y=Y0 + 10, max_h=110)
w2 = (CW - 18) / 2
B.kachel(s, CL, Y0 + 140, w2, 150, BLUE, "Warum drei Schichten und nicht eine",
         ["raw bleibt unverändert — der Vergleich mit dem Quellsystem bleibt möglich.",
          "stg ist als Sicht umgesetzt: immer frisch, kostet keinen Speicher.",
          "mart ist materialisiert — dort wird der Ladeschritt sichtbar, dort gehört später die Historisierung hin."])
B.kachel(s, CL + w2 + 18, Y0 + 140, w2, 150, METHOD, "Womit",
         ["Die Transformation ist SQL und gehört in versionierte Dateien, nicht in ein Werkzeugfenster.",
          "Der gesamte Bestand lädt in 2,3 Sekunden aus CSV und belegt danach 54,5 MB.",
          "Ein Data Warehouse würde hier Betrieb hinzufügen, ohne Fähigkeit hinzuzufügen."])

# ═══════════════════════════════════════════════════ Teil 4
kapitel("Ein Weg von vielen")

# Fünf Entscheidungen
s = neu("Slide")
B.kopf(s, "Rückblick", "Fünf Entscheidungen — und was jede von ihnen ausgeschlossen hat", Q_EIG)
B.zuordnung(s, Y0, [
    ("Wo auswerten?",
     "Gewählt: ein eigener Abzug in einer zweiten Datenbank. Verworfen: direkt im "
     "Kassensystem auswerten — Jahresabfragen hätten den laufenden Betrieb ausgebremst."),
    ("Wie modellieren?",
     "Gewählt: Galaxy-Schema mit zwei Granularitäten. Verworfen: One Big Table und "
     "Data Vault 2.0 — beide lösen andere Probleme als dieses."),
    ("Wie transformieren?",
     "Gewählt: SQL in versionierten Dateien, drei Schichten. Verworfen: eine Klickstrecke "
     "im ETL-Werkzeug — nicht vergleichbar, nicht wiederholbar."),
    ("Womit rechnen?",
     "Gewählt: DuckDB, ein Prozess, eine Datei. Verworfen: Data Warehouse, Data Lake, "
     "Lakehouse — bei 54,5 MB Bestand alles Aufwand ohne Gegenwert."),
    ("Wer verantwortet?",
     "Gewählt: ein zentrales Modell, eine Zuständigkeit. Verworfen: Data Mesh mit "
     "Datenprodukten je Domäne — es gibt hier nur eine Domäne."),
], rh=52.0)

# Drei Wege zu modellieren
s = neu("Lisa_Slide")
B.kopf(s, "Alternative 1 · Modellierung", "Dieselben Fakten lassen sich auf drei Arten ordnen", QUELLE)
einl(s, "Das Galaxy-Schema ist eine Wahl, keine Notwendigkeit. One Big Table spart jede "
        "Verknüpfung, kostet dafür Speicher und schließt die Positionsebene aus. Data Vault "
        "2.0 trennt Schlüssel, Beziehungen und Merkmale und hält jede Änderung historisiert — "
        "um den Preis, dass die Rohschicht ohne Hilfsschichten kaum abfragbar ist.")
bild(s, "09_modellierungswege", y=Y0, max_h=238)
B.band(s, Y0 + 250, 66, [
    "Der Bestand liegt in beiden Formen vor: obt_orders.csv hat 41 Spalten und 754.513 Zeilen "
    "und beantwortet die Filialfrage in 3,2 statt 3,9 Millisekunden. Der Unterschied ist "
    "messbar und bei dieser Größe bedeutungslos — die Positionsebene fehlt dafür ganz."])

# Warehouse und Lake
s = neu("Lisa_Slide")
B.kopf(s, "Alternative 2 · Architektur", "Zwei Erblinien: erst das Warehouse, dann der Lake", Q_LIT)
einl(s, "Beide Bauarten beantworten dieselbe Frage — wo die Auswertungsdaten liegen sollen — "
        "und beide scheitern an je einer Seite. Das Warehouse verlangt die Struktur vor dem "
        "Schreiben, der Lake verschiebt sie auf das Lesen. Das ist der ganze Gegensatz, alles "
        "Weitere folgt daraus.")
B.gegenueber(s, Y0, 190,
             (BLUE, "Data Warehouse · schema-on-write",
              ["Die Struktur steht fest, bevor eine Zeile geschrieben wird.",
               "Starke Zusicherungen: Transaktionen, Typen, Rechte, feste Kennzahlendefinitionen.",
               "Preis: unstrukturierte Daten passen nicht hinein, Speicher und Rechenwerk "
               "sind aneinander gebunden, ein Anbieterwechsel ist teuer."]),
             (HINT, "Data Lake · schema-on-read",
              ["Alles darf hinein, die Struktur entsteht erst beim Lesen.",
               "Billiger Objektspeicher, beliebige Rechenwerke, auch Bilder und Text.",
               "Preis: ohne Katalog und Zusicherungen entsteht ein Data Swamp — Dateien, "
               "die niemand mehr deuten kann."]),
             badge_l="1", badge_r="2")

# Lakehouse
s = neu("Lisa_Slide")
B.kopf(s, "Alternative 3 · Architektur", "Das Lakehouse legt Tabellensemantik über die Dateien", Q_LIT)
einl(s, "Das Lakehouse — 2021 von Armbrust, Ghodsi, Xin und Zaharia so benannt — behält den "
        "billigen Objektspeicher und holt die fehlenden Zusicherungen als eigene Schicht "
        "zurück: ein Tabellenformat über den Parquet-Dateien, dazu ein Katalog. Genau diese "
        "zwei Kästen unterscheiden es vom Lake.")
bild(s, "10_lakehouse_schichten", y=Y0, max_h=300)

# Lakehouse: Stand 2026
s = neu("Slide")
B.kopf(s, "Alternative 3 · Stand 2026", "Das Format ist entschieden, der Katalog ist es nicht", Q_LIT)
B.kachel(s, CL, Y0, w2, 150, BLUE, "Wo sich der Markt geeinigt hat",
         ["Apache Iceberg gilt als die Richtung, auf die sich die großen Anbieter zubewegen.",
          "Delta UniForm schreibt eine Tabelle einmal und legt Iceberg-Metadaten daneben.",
          "Über die Hälfte der Organisationen setzt nach Marktberichten Lakehouse-Muster ein."])
B.kachel(s, CL + w2 + 18, Y0, w2, 150, HINT, "Wo es offen ist",
         ["Der Katalog ist die neue Bindungsstelle: Polaris, Glue, Nessie, Unity, Lakekeeper.",
          "Apache Polaris ist seit Februar 2026 ein Projekt oberster Ebene der Apache-Stiftung.",
          "Wer heute wählt, wählt vor allem einen Katalog — nicht ein Dateiformat."])
B.band(s, Y0 + 170, 76, [
    "Für die Lehre ist das die eigentliche Beobachtung: Der Streit hat sich von der Frage "
    "nach dem Dateiformat zur Frage nach dem Verzeichnis verschoben. Offene Dateien allein "
    "schaffen keine Unabhängigkeit, solange der Katalog geschlossen bleibt — und am Katalog "
    "hängen Rechte, Schnappschüsse und der Schemaverlauf."])

# DuckLake
s = neu("Lisa_Slide")
B.kopf(s, "Alternative 4 · Gegenentwurf", "DuckLake stellt die Frage anders: Warum Metadaten in Dateien?", Q_DUCK)
einl(s, "Iceberg und Delta legen ihre Metadaten selbst als Dateien ab und brauchen für einen "
        "Lesevorgang mehrere Anfragen hintereinander. Für die Konsistenz mussten beide am Ende "
        "doch eine Datenbank als Katalog vorsehen. DuckLake zieht daraus den Schluss, alle "
        "Metadaten gleich in diese Datenbank zu legen — Version 1.0 seit April 2026.")
bild(s, "11_ducklake_metadaten", y=Y0, max_h=HOEHE)

# Data Mesh
s = neu("Lisa_Slide")
B.kopf(s, "Alternative 5 · Organisation", "Data Mesh ist kein Bauplan, sondern ein Organigramm", Q_MESH)
einl(s, "Data Mesh beantwortet nicht, wo die Daten liegen, sondern wer sie verantwortet: nicht "
        "ein zentrales Datenteam, sondern die Fachdomäne, die sie erzeugt — als Datenprodukt "
        "mit Zusagen, unter gemeinsamen Regeln auf einer geteilten Plattform. Das ist eine "
        "Aussage über Zuständigkeiten, nicht über Dateiformate.")
bild(s, "12_mesh", y=Y0, max_h=200)
B.klammer(s, CL, Y0 + 216, w2, 106, "Was davon geblieben ist",
          _paras(["Verantwortung liegt bei der Domäne, in der die Daten entstehen.",
                  "Daten gelten als Produkt mit Zusagen an ihre Abnehmer."]), GOOD)
B.klammer(s, CL + w2 + 18, Y0 + 216, w2, 106, "Was sich nicht durchgesetzt hat",
          _paras(["Nach Beraterberichten hat nur etwa ein Fünftel der Organisationen die nötige Reife.",
                  "In einer Erhebung von Oktober 2025 erreichten hybride Formen ihre Ziele häufiger."]), WARN)

# Bewertung
s = neu("Slide")
B.kopf(s, "Einordnung", "Jede Bauart kauft Tragfähigkeit mit Aufwand", Q_EIG)
bewertung(s, Y0,
          ["Datenmenge, die sie trägt", "Betriebsaufwand",
           "Bindung an einen Anbieter", "Aufwand für die Organisation"],
          [("Ein Prozess · DuckDB, DuckLake", [0.25, 0.10, 0.10, 0.10], True),
           ("Data Warehouse", [0.85, 0.60, 0.90, 0.40], False),
           ("Data Lake", [1.00, 0.75, 0.25, 0.40], False),
           ("Lakehouse", [1.00, 0.70, 0.30, 0.50], False),
           ("Data Mesh", [1.00, 0.90, 0.30, 1.00], False)])
B.band(s, Y0 + 286, 36, [
    "Längerer Balken = mehr davon. Die hervorgehobene Zeile ist der hier gewählte Weg."])

# Schluss
s = neu("Lisa_Slide")
B.kopf(s, "Ergebnis", "Warum hier ein einziger Prozess genügt — und woran man das Ende merkt", Q_MESS)
einl(s, "Die Wahl gegen jede dieser Architekturen ist keine grundsätzliche, sondern eine über "
        "die Größenordnung. Der gesamte Bestand — 754.513 Bestellungen und 2.950.082 "
        "Positionen — belegt nach dem Laden 54,5 MB und antwortet in Millisekunden. Ein "
        "verteiltes System hätte hier nichts zu verteilen.")
B.kennzahlen(s, Y0, 86, [
    ("CSV-Rohdaten", "311,4 MB", B.SEC),
    ("nach dem Laden", "54,5 MB", BLUE),
    ("Stern-Schema, 2 Verknüpfungen", "3,9 ms", BLUE),
    ("Positionsebene, 3 Verknüpfungen", "14,7 ms", BLUE),
])
B.kachel(s, CL, Y0 + 102, w2, 168, GOOD, "Warum das trägt",
         ["Analytische Arbeitsmengen sind meist klein: Der Median liegt nach Erhebungen "
          "unter 100 GB, dieser Bestand um drei Größenordnungen darunter.",
          "Ein Prozess, eine Datei, keine Zuständigkeitsgrenze — nichts davon muss betrieben werden.",
          "Parquet und SQL bleiben lesbar, auch wenn das Werkzeug gewechselt wird."])
B.kachel(s, CL + w2 + 18, Y0 + 102, w2, 168, WARN, "Woran man den Wechsel merkt",
         ["Der Bestand passt nicht mehr auf eine Maschine — dann Objektspeicher und Lakehouse.",
          "Viele schreiben gleichzeitig — dann ein Tabellenformat mit Transaktionen.",
          "Mehrere Fachbereiche streiten über Kennzahlendefinitionen — dann ist es eine "
          "Organisationsfrage, und erst dann lohnt der Blick auf Data Mesh."])

for sl in prs.slides:
    aufraeumen(sl)
prs.save(str(OUT))
print(f"Geschrieben: {OUT.name}")
print(f"Folien: {len(prs.slides._sldIdLst)}")
if B.WARNINGS:
    print("\nWarnungen der Bausteine:")
    for x in B.WARNINGS:
        print("  -", x)
else:
    print("Keine Baustein-Warnungen.")
