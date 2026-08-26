#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
deckwerk.py — gemeinsames Geruest der Bauskripte in diesem Ordner.

Hier steht, was bau_datenmodell.py und bau_bi.py beide brauchen: den Skill
finden, den Platzhalterfehler des spec_loader umgehen, Absaetze richtig
uebergeben, Diagramme und Bildschirmfotos massstabsgetreu setzen. Ohne diese
Datei stuenden dieselben vierzig Zeilen zweimal da und wuerden auseinanderlaufen.

Verwendung:

    import deckwerk as D
    deck = D.Deck()
    s = deck.neu("Lisa_Slide")
    D.B.kopf(s, "Kicker", "Aussagentitel", D.QUELLE)
    deck.einl(s, "Drei bis vier Zeilen ...")
    deck.bild(s, "07_galaxy", max_h=D.HOEHE)
    deck.speichern(pfad)
"""
import importlib.util
import os
import struct
import sys
import types
from pathlib import Path

HIER = Path(__file__).resolve().parent
DIA = HIER / "diagramme"
BILDER = HIER / "bilder"


def finde_skill():
    if os.environ.get("THWS_SKILL"):
        p = Path(os.environ["THWS_SKILL"])
        if (p / "assets" / "template.pptx").exists():
            return p
        sys.exit(f"FEHLER: THWS_SKILL zeigt auf {p}, dort fehlt assets/template.pptx")
    for basis in (Path.home() / ".claude" / "skills",
                  Path.home() / "Library/Application Support/Claude"):
        treffer = sorted(basis.glob("**/thws-slides/assets/template.pptx"))
        if treffer:
            return treffer[-1].parent.parent
    sys.exit('FEHLER: Skill thws-slides nicht gefunden. '
             'export THWS_SKILL="/Pfad/zu/skills/thws-slides"')


SKILL = finde_skill()
SPECS = SKILL / "assets" / "specs"
TEMPLATE = SKILL / "assets" / "template.pptx"


def _load_spec(name="bint"):
    sp = importlib.util.spec_from_file_location(f"spec_{name}", SPECS / f"{name}.py")
    m = importlib.util.module_from_spec(sp)
    sp.loader.exec_module(m)
    return m


# Der spec_loader des Skills zeigt auf <skill>/specs, die Specs liegen aber
# unter <skill>/assets/specs. Ohne diesen Ersatz enden alle Bausteine mit
# "Spec 'bint' nicht gefunden".
_fake = types.ModuleType("spec_loader")
_fake.load = _load_spec
_fake.SPECS_DIR = SPECS
sys.modules["spec_loader"] = _fake

sys.path.insert(0, str(SKILL / "scripts"))
import bausteine as B                                    # noqa: E402
from pptx import Presentation                            # noqa: E402
from pptx.util import Pt                                 # noqa: E402
from pptx.enum.text import PP_ALIGN                      # noqa: E402

G = B.G
BLUE, WARN, GOOD, METHOD, HINT = G.BLUE, G.WARN, G.GOOD, G.METHOD, G.HINT
SEC, WEISS, HAAR = B.SEC, B.WHITE, B.HAIR
CL, CW, CR = G.CONTENT_LEFT, G.CONTENT_WIDTH, G.CONTENT_RIGHT
Y0, YMAX = G.CONTENT_ZONE_Y_MIN, G.CONTENT_ZONE_Y_MAX
HOEHE = YMAX - Y0

QUELLE = "Eigene Darstellung · Diagramme aus slides/diagramme/*.mmd"
Q_FOTO = "Bildschirmfoto aus dem Projekt · aufgenommen mit slides/screenshots.mjs"
Q_MESS = "Eigene Messung · DuckDB 1.5.5 auf dataset/, Median aus sieben Läufen"


# --- Absatzformat -----------------------------------------------------------
# as_paras() der Bausteine erkennt nur Listen von LISTEN als mehrere Absaetze.
# Flache Stringlisten liefen sonst zu einem einzigen Absatz zusammen.
def paras(x):
    if isinstance(x, list) and x and all(isinstance(e, str) for e in x):
        return [[e] for e in x]
    return x


_kachel, _band, _gegenueber = B.kachel, B.band, B.gegenueber
B.kachel = lambda s, x, y, w, h, f, t, b, gap=5: _kachel(s, x, y, w, h, f, t, paras(b), gap)
B.band = lambda s, y, h, p, **k: _band(s, y, h, paras(p), **k)
B.gegenueber = lambda s, y, h, l, r, **k: _gegenueber(
    s, y, h, (l[0], l[1], paras(l[2])), (r[0], r[1], paras(r[2])), **k)


def png_groesse(p):
    """Breite und Hoehe aus dem Dateikopf — ohne Bildbibliothek."""
    d = open(p, "rb").read(64)
    if d[:8] == b"\x89PNG\r\n\x1a\n":
        return struct.unpack(">II", d[16:24])
    # JPEG: bis zum ersten SOF-Marker vorspulen
    with open(p, "rb") as f:
        f.seek(2)
        while True:
            b0 = f.read(1)
            while b0 and b0 != b"\xff":
                b0 = f.read(1)
            m = f.read(1)
            while m == b"\xff":
                m = f.read(1)
            if m and 0xC0 <= m[0] <= 0xCF and m[0] not in (0xC4, 0xC8, 0xCC):
                f.read(3)
                h, w = struct.unpack(">HH", f.read(4))
                return w, h
            laenge = struct.unpack(">H", f.read(2))[0]
            f.seek(laenge - 2, 1)


class Deck:
    def __init__(self):
        self.prs = Presentation(str(TEMPLATE))
        self.layout = {l.name: l for l in self.prs.slide_layouts}

    # --- Folien ---------------------------------------------------------
    def neu(self, name):
        return self.prs.slides.add_slide(self.layout[name])

    def kapitel(self, titel, groesse=29):
        s = self.neu("Chapter")
        for sh in s.placeholders:
            if str(sh.placeholder_format.type).startswith("TITLE"):
                sh.text_frame.text = titel
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(groesse)
        self.aufraeumen(s)
        return s

    def deckblatt(self, oben, unten):
        s = self.neu("Frontpage_Digital")
        for sh in s.placeholders:
            i = sh.placeholder_format.idx
            if i == 10:
                sh.text_frame.text = oben
            elif i == 11:
                sh.text_frame.text = unten
        self.aufraeumen(s)
        return s

    def einl(self, slide, text, variante="Lisa_Slide"):
        """Einleitung in den Master-Platzhalter idx 14, auf 14 pt gesetzt.

        Die Geometrie wird vollstaendig gesetzt: Wird nur die Hoehe geaendert,
        entsteht ein xfrm ohne Position und der Platzhalter verschwindet beim
        Rendern.
        """
        ph = B.einleitung(slide, [text], variante)
        vor = [q for q in self.layout[variante].placeholders
               if q.placeholder_format.idx == 14][0]
        ph.left, ph.top, ph.width = vor.left, vor.top, vor.width
        ph.height = Pt(63.8)
        for p in ph.text_frame.paragraphs:
            p.alignment = PP_ALIGN.JUSTIFY
            for r in p.runs:
                r.font.size = Pt(14)
        return ph

    @staticmethod
    def aufraeumen(slide):
        for sh in list(slide.placeholders):
            try:
                if not sh.text_frame.text.strip():
                    sh._element.getparent().remove(sh._element)
            except Exception:
                pass

    # --- Bilder ---------------------------------------------------------
    def _setze(self, slide, pfad, y, max_h, x, max_w, rahmen):
        pw, ph = png_groesse(pfad)
        y = Y0 if y is None else y
        max_h = (YMAX - y) if max_h is None else max_h
        max_w = CW if max_w is None else max_w
        sk = min(max_w / pw, max_h / ph)
        w, h = pw * sk, ph * sk
        x = CL if x is None else x
        if rahmen:
            B.shape(slide, B.MSO_SHAPE.RECTANGLE, x - 1, y - 1, w + 2, h + 2, None, HAAR)
        slide.shapes.add_picture(str(pfad), B.E(x), B.E(y), B.E(w), B.E(h))
        return w, h, sk

    def bild(self, slide, name, y=None, max_h=None, x=None, max_w=None):
        """Mermaid-Diagramm aus diagramme/.

        Der Schriftgrad im Bild folgt aus dem Massstab: Mermaid setzt 15 CSS-px,
        die PNG sind dreifach ueberabgetastet. Auf der Folie kommen davon
        45 * Massstab Punkte an — unter 11 pt ist das im Hoersaal nicht lesbar.
        """
        p = DIA / f"{name}.png"
        if not p.exists():
            sys.exit(f"FEHLER: {p.name} fehlt — erst 'node render_mermaid.mjs' ausfuehren.")
        w, h, sk = self._setze(slide, p, y, max_h, x, max_w, False)
        if 45 * sk < 10.5:
            B.WARNINGS.append(f"Diagramm {name}: Beschriftung nur {45 * sk:.1f} pt")
        return w, h

    def foto(self, slide, name, y=None, max_h=None, x=None, max_w=None):
        """Bildschirmfoto aus bilder/, mit Haarlinie umrandet.

        Der Rahmen ist noetig, weil die Aufnahmen an den Raendern hell sind und
        sonst in die Folie ausfransen.
        """
        treffer = [p for p in (BILDER / f"{name}.png", BILDER / f"{name}.jpg") if p.exists()]
        if not treffer:
            sys.exit(f"FEHLER: bilder/{name}.png|.jpg fehlt — erst 'node screenshots.mjs' ausfuehren.")
        w, h, _ = self._setze(slide, treffer[0], y, max_h, x, max_w, True)
        return w, h

    # --- eigene Motive --------------------------------------------------
    @staticmethod
    def bewertung(slide, y, spalten, zeilen, rh=44.0, bw=250.0):
        """Bewertungsmatrix: Bezeichner links, Balken je Spalte.

        Balken statt Harvey-Kreisen: Ein teilgefuellter Kreis ist eine PIE-Form,
        deren Winkel LibreOffice beim Rendern verwirft — halbvolle Kreise kamen
        dort als leere oder volle heraus. Rechtecke stellt jedes Programm gleich dar.

        zeilen: Liste von (Bezeichner, [Anteil 0..1 je Spalte], Hervorhebung).
        """
        n = len(spalten)
        sw = (CR - CL - bw) / n
        balken_w = min(100.0, sw - 34)
        for j, kopfzeile in enumerate(spalten):
            B.textbox(slide, CL + bw + j * sw, y, sw, 30, [(kopfzeile, False, SEC)],
                      10.5, align=PP_ALIGN.CENTER)
        yy = y + 34
        for bez, werte, hl in zeilen:
            farbe = BLUE if hl else SEC
            B.shape(slide, B.MSO_SHAPE.RECTANGLE, CL, yy, bw + n * sw, rh,
                    "F4F7FA" if hl else WEISS, HAAR)
            B.textbox(slide, CL + 14, yy + 6, bw - 20, rh - 12, [(bez, hl, farbe)],
                      12, anchor=B.MSO_ANCHOR.MIDDLE)
            for j, a in enumerate(werte):
                bx = CL + bw + j * sw + (sw - balken_w) / 2
                by = yy + rh / 2 - 4
                B.shape(slide, B.MSO_SHAPE.RECTANGLE, bx, by, balken_w, 8, "E8E6E0")
                if a > 0:
                    B.shape(slide, B.MSO_SHAPE.RECTANGLE, bx, by, balken_w * a, 8, farbe)
            yy += rh + 6
        return yy

    # --- Abschluss ------------------------------------------------------
    def speichern(self, ziel):
        for sl in self.prs.slides:
            self.aufraeumen(sl)
        self.prs.save(str(ziel))
        print(f"Geschrieben: {Path(ziel).name}")
        print(f"Folien: {len(self.prs.slides._sldIdLst)}")
        if B.WARNINGS:
            print("\nWarnungen der Bausteine:")
            for x in B.WARNINGS:
                print("  -", x)
            return 1
        print("Keine Baustein-Warnungen.")
        return 0
