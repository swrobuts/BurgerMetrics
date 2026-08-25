#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
bau_deck.py — erzeugt das BurgerMetrics-Foliendeck nach dem THWS-Folienmaster.

    python3 bau_deck.py                      # -> ../../BurgerMetrics_Fallstudie.pptx
    python3 bau_deck.py -o /pfad/deck.pptx

Voraussetzungen
---------------
    pip install python-pptx
    Skill "thws-slides" mit assets/template.pptx und scripts/bausteine.py

Der Skill-Pfad kommt aus der Umgebungsvariablen THWS_SKILL; ohne sie wird die
uebliche Stelle der lokalen Installation durchsucht. Vorlage und Hausschrift
liegen bewusst ausserhalb dieses Repositorys — siehe README.md in diesem Ordner.

Bekannter Fehler im Skill
-------------------------
spec_loader.SPECS_DIR zeigt auf <skill>/specs, die Specs liegen aber unter
<skill>/assets/specs. Statt die Installation zu aendern, haengt dieses Skript
einen korrigierten Loader in sys.modules ein, bevor bausteine importiert wird.
"""
import argparse
import importlib.util
import os
import sys
import types
from pathlib import Path


def finde_skill():
    if os.environ.get("THWS_SKILL"):
        p = Path(os.environ["THWS_SKILL"])
        if (p / "assets" / "template.pptx").exists():
            return p
        sys.exit(f"FEHLER: THWS_SKILL zeigt auf {p}, dort fehlt assets/template.pptx")
    basis = Path.home() / "Library/Application Support/Claude"
    treffer = sorted(basis.glob("**/skills/thws-slides/assets/template.pptx"))
    if treffer:
        return treffer[-1].parent.parent
    sys.exit(
        "FEHLER: Der Skill thws-slides wurde nicht gefunden.\n"
        "Pfad setzen, zum Beispiel:\n"
        '  export THWS_SKILL="/Pfad/zu/skills/thws-slides"'
    )


SKILL = finde_skill()
SPECS = SKILL / "assets" / "specs"


def _load_spec(name="bint"):
    sp = importlib.util.spec_from_file_location(f"spec_{name}", SPECS / f"{name}.py")
    m = importlib.util.module_from_spec(sp)
    sp.loader.exec_module(m)
    return m


_fake = types.ModuleType("spec_loader")
_fake.load = _load_spec
_fake.SPECS_DIR = SPECS
sys.modules["spec_loader"] = _fake

sys.path.insert(0, str(SKILL / "scripts"))
import bausteine as B                                    # noqa: E402
from pptx import Presentation                            # noqa: E402
from pptx.util import Pt, Emu                            # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                   # noqa: E402
from pptx.enum.text import PP_ALIGN                      # noqa: E402

_ap = argparse.ArgumentParser(description="BurgerMetrics-Foliendeck bauen")
_ap.add_argument("-o", "--ausgabe", default=None,
                 help="Zieldatei (Standard: ../../BurgerMetrics_Fallstudie.pptx)")
_args = _ap.parse_args()

G = B.G


# --- Absatzformat richtigstellen -------------------------------------------
# as_paras() erkennt nur eine Liste von LISTEN als mehrere Absaetze. Eine
# flache Stringliste wird zu einem einzigen Absatz mit mehreren Runs — die
# Aufzaehlungspunkte liefen dadurch ohne Trennung ineinander. Die Bausteine
# werden deshalb so umhuellt, dass sie flache Listen korrekt aufteilen.
def _paras(x):
    if isinstance(x, list) and x and all(isinstance(e, str) for e in x):
        return [[e] for e in x]
    return x


_kachel, _klammer, _gegenueber, _band = B.kachel, B.klammer, B.gegenueber, B.band


def _k(slide, x, y, w, h, farbe, titel, body, gap=5):
    return _kachel(slide, x, y, w, h, farbe, titel, _paras(body), gap)


def _kl(slide, x, y, w, h, titel, punkte, farbe=None):
    return _klammer(slide, x, y, w, h, titel, _paras(punkte),
                    farbe if farbe is not None else B.BLUE)


def _g(slide, y, h, links, rechts, **kw):
    return _gegenueber(slide, y, h,
                       (links[0], links[1], _paras(links[2])),
                       (rechts[0], rechts[1], _paras(rechts[2])), **kw)


def _b(slide, y, h, paras, **kw):
    return _band(slide, y, h, _paras(paras), **kw)


B.kachel, B.klammer, B.gegenueber, B.band = _k, _kl, _g, _b
# ---------------------------------------------------------------------------
BLUE, WARN, GOOD, METHOD, HINT = G.BLUE, G.WARN, G.GOOD, G.METHOD, G.HINT
CL, CW = G.CONTENT_LEFT, G.CONTENT_WIDTH
CR = G.CONTENT_RIGHT
Y0 = G.CONTENT_ZONE_Y_MIN          # 172.0 laut Spec (Skill-Altwert 176)
YMAX = G.CONTENT_ZONE_Y_MAX

TEMPLATE = SKILL / "assets" / "template.pptx"
OUT = Path(_args.ausgabe) if _args.ausgabe else (
    Path(__file__).resolve().parent.parent.parent / "BurgerMetrics_Fallstudie.pptx")

QUELLE = "Eigene Darstellung · Datenstand BurgerMetrics, 25.08.2026"

prs = Presentation(str(TEMPLATE))
LAYOUT = {l.name: l for l in prs.slide_layouts}


def neu(layout_name):
    return prs.slides.add_slide(LAYOUT[layout_name])


def einl(slide, paras, variante="Lisa_Slide"):
    """Einleitung setzen, auf 14 pt, mit vollständiger Geometrie.

    Wird nur die Höhe gesetzt, entsteht ein xfrm ohne Position und der
    Platzhalter wird beim Rendern nicht gezeichnet (folienmaster.md).
    """
    ph = B.einleitung(slide, paras, variante)
    vor = [q for q in LAYOUT[variante].placeholders
           if q.placeholder_format.idx == 14][0]
    ph.left, ph.top, ph.width = vor.left, vor.top, vor.width
    ph.height = Pt(63.8)                      # vier Zeilen
    for p in ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.JUSTIFY
        for r in p.runs:
            r.font.size = Pt(14)
    return ph


def aufraeumen(slide):
    """Leere Platzhalter entfernen — sie erscheinen sonst als Geisterzeilen."""
    for sh in list(slide.placeholders):
        try:
            if not sh.text_frame.text.strip():
                sh._element.getparent().remove(sh._element)
        except Exception:
            pass


def kapitel(titel):
    s = neu("Chapter")
    for sh in s.placeholders:
        if str(sh.placeholder_format.type).startswith("TITLE"):
            sh.text_frame.text = titel
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(29)
    aufraeumen(s)
    return s


# ══════════════════════════════════════════════ Deckblatt
s = neu("Frontpage_Digital")
for sh in s.placeholders:
    f = sh.placeholder_format
    if f.idx == 10:
        sh.text_frame.text = "BurgerMetrics"
    elif f.idx == 11:
        sh.text_frame.text = "Von der Kasse bis zum Bericht"
aufraeumen(s)

# ══════════════════════════════════════════════ Kapitel 1
kapitel("Ein Lehrdatensatz wird entworfen, nicht gefunden")

# --- 1 Anforderungen (kachel)
s = neu("Lisa_Slide")
B.kopf(s, "Entwurf", "Groß genug für Modellierung, klein genug für Nachrechenbarkeit", QUELLE)
einl(s, ["Ein Lehrdatensatz muss zwei Bedingungen erfüllen, die einander widersprechen. "
         "Ist er klein und sauber, lässt er sich überblicken, ohne dass ein Datenmodell "
         "nötig wäre — die Einsicht, um die es geht, entsteht dann nicht. Sind es echte "
         "Unternehmensdaten, darf man sie nicht ausliefern. Beides zusammen geht nur synthetisch."])
w = (CW - 2 * 18) / 3
for i, (t, b, c) in enumerate([
    ("Groß genug", ["754.513 Bestellungen", "2.950.082 Positionen",
                     "neun Jahre, acht Filialen", "Auswertung ohne Modell nicht möglich"], BLUE),
    ("Nachrechenbar", ["Jede Zahl aus mitgelieferten Dateien",
                       "keine Blackbox, kein Fernzugriff", "79 Angaben maschinell geprüft"], GOOD),
    ("Synthetisch", ["kein Datenschutzproblem", "Muster gezielt eingebaut",
                     "Preis: sauberer als die Wirklichkeit"], HINT),
]):
    B.kachel(s, CL + i * (w + 18), Y0, w, 150, c, t, b)
B.band(s, Y0 + 168, 54, ["Die Auswahl der eingebauten Muster ist die eigentliche Entwurfsarbeit — "
                         "nicht die Datenmenge."])

# --- 2 Systemkontext (chevron_kette)
s = neu("Slide")
B.kopf(s, "Entwurf", "Drei Anwendungen bilden drei Rollen ab und speisen denselben Bestand", QUELLE)
B.chevron_kette(s, Y0 + 30, 66, [
    "POS-Terminal", "Online-Shop", "CSV-Bestand", "BI-Bericht"], farbe=BLUE)
w2 = (CW - 18) / 2
B.kachel(s, CL, Y0 + 130, w2, 142, METHOD, "Was die Anwendungen leisten",
         ["Das POS-Terminal zeigt bei jedem Bedienschritt, welche Datenzeile dabei entsteht.",
          "Der Shop bildet dieselbe Bestellung aus Kundensicht ab.",
          "Der Bericht wertet aus, was beide erzeugt haben."])
B.kachel(s, CL + w2 + 18, Y0 + 130, w2, 142, GOOD, "Warum das zusammengehört",
         ["Eine Kennzahl im Bericht soll bis zu dem Vorgang zurückverfolgbar sein, der sie erzeugt hat.",
          "Wer sieht, dass Drive-Through höhere Bestellwerte hat, kann an der Kasse nachsehen, wie eine solche Bestellung entsteht."])

# --- 3 Muster (kachel)
s = neu("Lisa_Slide")
B.kopf(s, "Entwurf", "Auch Scheinmuster sind mit Absicht eingebaut", QUELLE)
einl(s, ["Ein synthetischer Datensatz enthält genau die Muster, die man hineinlegt. Für die "
         "Lehre kommt es darauf an, nicht nur bestätigungsfähige Muster anzulegen: Wer immer "
         "findet, was er sucht, lernt das Prüfen nicht. Deshalb liegen drei Sorten nebeneinander."])
for i, (t, b, c) in enumerate([
    ("Bestätigen sich", ["Wochentagseffekte, Mittagspeak", "Eis im Sommer: Faktor 2,78",
                          "Bargeld fällt von 49 auf 18 Prozent"], GOOD),
    ("Nur mit Sorgfalt", ["Wettereffekt liegt unter dem Wachstumstrend",
                          "Uni-Standort: minus 31 Prozent im August",
                          "sichtbar nur im Vergleich mit sich selbst"], HINT),
    ("Halten nicht stand", ["Bier bei Burgern: Lift 1,03",
                            "Kundenherkunft gleichverteilt",
                            "App-Nutzung überall 43 bis 45 Prozent"], WARN),
]):
    B.kachel(s, CL + i * (w + 18), Y0, w, 150, c, t, b)

# ══════════════════════════════════════════════ Kapitel 2
kapitel("Vom operativen Modell zum Auswertungsmodell")

# --- 4 Reduktion (trichter)
s = neu("Slide")
B.kopf(s, "Modellierung", "Aus 26 operativen Tabellen werden zwölf", QUELLE)
B.trichter(s, Y0 + 20, 150, [
    "26 Tabellen im ERP-Modell (3NF), sieben fachliche Bereiche",
    "Absatzpfad, Stammdaten, externe Daten — die übrigen Bereiche beantworten keine der Leitfragen",
    "12 Tabellen: zwei Faktentabellen, zehn Dimensionen"], farbe=BLUE)
B.band(s, Y0 + 195, 62, [
    "Ein Auswertungsmodell bildet nicht den Betrieb ab, sondern die Fragen, die an ihn gestellt "
    "werden. Einkauf, Lager, Preishistorie und Schichtplan bleiben im Quellsystem — mit einer "
    "Nebenwirkung, die auf der übernächsten Folie wieder auftaucht."])

# --- 5 Denormalisierung (gegenueber)
s = neu("Lisa_Slide")
B.kopf(s, "Modellierung", "Denormalisierung kostet Redundanz und spart drei Verknüpfungen", QUELLE)
einl(s, ["In der dritten Normalform steht jedes Merkmal genau einmal — die Kategorie eines Artikels "
         "in ihrer eigenen Tabelle. Das ist beim Schreiben richtig und bei der Auswertung teuer. "
         "Vertretbar wird die Umkehrung dadurch, dass ein Auswertungsbestand periodisch neu beladen "
         "und nicht laufend fortgeschrieben wird."])
B.gegenueber(s, Y0, 156,
             (BLUE, "Operativ (3NF)",
              ["artikel, artikelkategorie und artikelunterkategorie getrennt",
               "Umsatz je Kategorie braucht vier Tabellen und drei Verknüpfungen",
               "Aenderung eines Kategorienamens trifft genau eine Zeile"]),
             (METHOD, "Analytisch (denormalisiert)",
              ["drei Tabellen werden zu dim_product",
               "dieselbe Frage braucht eine Verknüpfung",
               "Kategorie steht redundant in jeder der 57 Produktzeilen"]),
             badge_l="1", badge_r="2")

# --- 6 Galaxy-Schema (eigenes Diagramm)
s = neu("Slide")
B.kopf(s, "Modellierung", "Zwei Granularitäten verlangen zwei Faktentabellen", QUELLE)

DW, DH = 168, 34
FW, FH = 268, 46
fx = CL + (CW - FW) / 2                     # Fakten mittig
fy1, fy2 = Y0 + 60, Y0 + 178


def kasten(x, y, w, h, text, fill, tcol, size=11, bold=True, gestrichelt=False):
    sh = B.shape(s, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, G.HAIRLINE)
    if gestrichelt:
        sh.line.dash_style = 4               # MSO_LINE.DASH
    B.label(sh, [[(text, bold, tcol)]], size)
    return sh


def linie(x1, y1, x2, y2):
    ln = s.shapes.add_connector(1, Emu(int(x1 * 12700)), Emu(int(y1 * 12700)),
                                Emu(int(x2 * 12700)), Emu(int(y2 * 12700)))
    ln.line.color.rgb = B.rgb(G.HAIRLINE)
    ln.line.width = Pt(1)
    return ln


# drei Dimensionen oben, Linie senkrecht auf die Oberkante von fact_orders
oben = ["dim_date", "dim_weather", "dim_branch"]
ox = [CL + 40, CL + 40 + (DW + 30), CL + 40 + 2 * (DW + 30)]
for nm, x in zip(oben, ox):
    kasten(x, Y0, DW, DH, nm, "FFFFFF", G.TEXT_BODY, 11, False)
    linie(x + DW / 2, Y0 + DH, x + DW / 2, fy1 - 10)
linie(ox[0] + DW / 2, fy1 - 10, ox[2] + DW / 2, fy1 - 10)   # Sammelschiene
linie(fx + FW / 2, fy1 - 10, fx + FW / 2, fy1)

# drei Dimensionen rechts, waagerecht an die rechte Kante
rechts = ["dim_customer", "dim_payment_method", "dim_promotion"]
rx = CR - DW
for k, nm in enumerate(rechts):
    ry = fy1 - 26 + k * (DH + 12)
    kasten(rx, ry, DW, DH, nm, "FFFFFF", G.TEXT_BODY, 11, False)
    linie(rx, ry + DH / 2, rx - 26, ry + DH / 2)
linie(rx - 26, fy1 - 26 + DH / 2, rx - 26, fy1 - 26 + 2 * (DH + 12) + DH / 2)
linie(rx - 26, fy1 + FH / 2, fx + FW, fy1 + FH / 2)

# die beiden Faktentabellen
kasten(fx, fy1, FW, FH, "fact_orders · 754.513 Bestellungen", BLUE, "FFFFFF", 11)
kasten(fx, fy2, FW, FH, "fact_order_items · 2.950.082 Positionen", BLUE, "FFFFFF", 11)
linie(fx + FW / 2, fy1 + FH, fx + FW / 2, fy2)
B.textbox(s, fx + FW / 2 + 8, (fy1 + FH + fy2) / 2 - 8, 100, 14,
          [("order_id", False, HINT)], 10)

# dim_product an die Positionsebene
px = fx + FW + 40
kasten(px, fy2 + (FH - DH) / 2, DW, DH, "dim_product", "FFFFFF", G.TEXT_BODY, 11, False)
linie(fx + FW, fy2 + FH / 2, px, fy2 + FH / 2)

# die drei nicht angebundenen Dimensionen, gestrichelt abgesetzt
for k, (nm, hw) in enumerate([("dim_employee", "nur an dim_branch"),
                              ("dim_time_slot", "anschließbar, ungenutzt"),
                              ("dim_supplier", "Orphan Dimension")]):
    x = CL + 40 + k * (DW + 30)
    sh = B.shape(s, MSO_SHAPE.RECTANGLE, x, Y0 + 262, DW, 34, "FFFFFF", G.HAIRLINE)
    sh.line.dash_style = 4
    B.label(sh, [[(nm, True, G.TEXT_BODY)], [(hw, False, WARN)]], 9.5)

# --- 7 Schiefe Dimensionen (kachel)
s = neu("Lisa_Slide")
B.kopf(s, "Modellierung", "Drei Dimensionen hängen bewusst schief", QUELLE)
einl(s, ["Nicht jede Unsauberkeit im Modell ist ein Versäumnis. Drei Dimensionen sind absichtlich "
         "nicht sauber angebunden, weil sich an ihnen zeigen lässt, wie solche Zustände entstehen "
         "und was sie für die Beantwortbarkeit von Fragen bedeuten. Sie sind Diskussionsmaterial, "
         "kein Fehler."])
for i, (t, b, c) in enumerate([
    ("dim_supplier", ["Orphan Dimension — kein Fremdschlüssel, nirgends",
                      "im operativen Modell gibt es die Beziehung über den Einkaufspfad",
                      "sie ging beim Uebergang verloren"], WARN),
    ("dim_employee", ["hängt nur an dim_branch",
                      "keine employee_id in fact_orders",
                      "Umsatz je Mitarbeiter ist damit nicht beantwortbar"], HINT),
    ("dim_time_slot", ["über hour anschließbar",
                       "wird im Bestand nirgends genutzt",
                       "die One Big Table führt hour, nicht time_slot"], METHOD),
]):
    B.kachel(s, CL + i * (w + 18), Y0, w, 150, c, t, b)

# ══════════════════════════════════════════════ Kapitel 3
kapitel("Auswerten — und die Fallen, die dabei aufgehen")

# --- 8 Fan Trap (kennzahlen + band)
s = neu("Slide")
B.kopf(s, "Auswertung", "Wer die Ebenen vermischt, erhält das 4,88-fache — ohne Fehlermeldung", QUELLE)
B.kennzahlen(s, Y0, 104, [
    ("richtig gerechnet (€)", "14.522.378,70", GOOD),
    ("nach dem Fan Trap (€)", "70.809.660,60", WARN),
    ("Verzerrungsfaktor", "4,88", WARN),
    ("Positionen je Bestellung", "3,91", BLUE),
])
B.kachel(s, CL, Y0 + 124, w2, 136, WARN, "Die Abfrage, die harmlos aussieht",
         ["SELECT SUM(o.net_total) FROM fact_orders o JOIN fact_order_items i ON o.order_id = i.order_id;",
          "Der Join vervielfacht jede Bestellung entlang ihrer Positionen. Es gibt keine Fehlermeldung."])
B.kachel(s, CL + w2 + 18, Y0 + 124, w2, 136, GOOD, "Warum der Faktor nicht 3,91 ist",
         ["Bestellungen mit vielen Positionen haben im Mittel auch höhere Beträge und schlagen beim Vervielfachen stärker durch.",
          "Der Fehler lässt sich deshalb nicht durch Teilen zurückrechnen — die Abfrage muss korrigiert werden."])

# --- 9 Würfeloperationen (zuordnung)
s = neu("Lisa_Slide")
B.kopf(s, "Auswertung", "Die sechs Würfeloperationen sind GROUP BY, WHERE und ein Ebenenwechsel", QUELLE)
einl(s, ["Die Operationen am OLAP-Würfel wirken wie eine eigene Begriffswelt neben SQL. Tatsächlich "
         "lassen sie sich vollständig auf drei Bausteine zurückführen. Wer das einmal gesehen hat, "
         "übersetzt zwischen Werkzeugoberfläche und Abfrage, ohne beides getrennt lernen zu müssen."])
B.zuordnung(s, Y0, [
    ("Roll-up", "gröberes GROUP BY — vom Quartal auf das Jahr"),
    ("Drill-Down", "feineres GROUP BY — vom Jahr auf das Quartal"),
    ("Slice", "WHERE auf genau eine Dimension"),
    ("Dice", "WHERE auf mehrere Dimensionen gleichzeitig"),
    ("Pivot", "dieselben Daten, vertauschte Achsen (CASE-Aggregation)"),
    ("Drill-Through", "Aggregation verlassen — zurück zum Einzelbeleg"),
], rh=38.0, gap=6.0)

# --- 10 Lift (vier Kacheln, je eine Regel)
s = neu("Lisa_Slide")
B.kopf(s, "Auswertung", "Bier zum Burger sieht nach Muster aus — der Lift sagt das Gegenteil", QUELLE)
einl(s, ["Support sagt, wie oft zwei Produkte gemeinsam vorkommen, Konfidenz, wie oft das eine auf "
         "das andere folgt. Beide können hoch sein, ohne dass ein Zusammenhang besteht. Erst der "
         "Lift setzt die Konfidenz ins Verhältnis zur Grundrate: Bei 1,0 sind die Produkte "
         "unabhängig, darunter schließen sie einander eher aus."])
wl = (CW - 3 * 14) / 4
for i2, (t, b, c) in enumerate([
    ("Burger + Fries + Cola", ["Support 29,2 %", "Lift 2,30",
                               "der erwartete Klassiker"], GOOD),
    ("Nuggets → BBQ-Sauce", ["Support 2,7 %", "Konfidenz 40,9 %", "Lift 8,93",
                             "kleinster Support, stärkster Zusammenhang"], GOOD),
    ("Breakfast → Coffee", ["Support 4,1 %", "Konfidenz 60,8 %", "Lift 3,59",
                            "hohe Konfidenz, echter Effekt"], GOOD),
    ("Burger → Bier", ["Support 7,6 %", "Konfidenz 12,7 %", "Lift 1,03",
                       "Grundrate Bier: 12,3 % — kein Effekt"], WARN),
]):
    B.kachel(s, CL + i2 * (wl + 14), Y0, wl, 158, c, t, b)
B.band(s, Y0 + 176, 54, [
    "Wer nur auf Support und Konfidenz schaut, findet in der letzten Kachel eine Regel, die keine "
    "ist — und übersieht in der zweiten die stärkste des Datensatzes."])

# --- 11 Basiseffekt (gegenueber)
s = neu("Lisa_Slide")
B.kopf(s, "Auswertung", "Gestaffelte Eröffnungen erzeugen Wachstumsraten ohne Wachstum", QUELLE)
einl(s, ["Die acht Filialen eröffnen zwischen März 2017 und März 2023. Ein Vorjahresvergleich "
         "trifft deshalb bei jeder Filiale einmal auf ein Rumpfjahr. Die Veränderungsrate misst dann "
         "vor allem die Länge des Vergleichszeitraums — ein Basiseffekt, kein Geschäftsverlauf."])
B.gegenueber(s, Y0, 156,
             (WARN, "Was die Zahl behauptet",
              ["BM Zellerau: über 20 Prozent Zuwachs von 2023 auf 2024",
               "absolute Filialumsätze zwischen 0,75 und 3,39 Millionen Euro",
               "gelesen als Leistungsunterschied der Standorte"]),
             (GOOD, "Was tatsächlich gemessen wird",
              ["2023 umfasst nur zehn Monate — Rumpfjahr gegen volles Jahr",
               "Europastern läuft neun Jahre, Zellerau drei",
               "ohne Normierung misst der Vergleich vor allem die Betriebsdauer"]),
             badge_l="!", badge_r="✓")

# ══════════════════════════════════════════════ Kapitel 4
kapitel("Prüfen ist Handwerk, nicht Vertrauen")

# --- 12 Drei Rechenwege (zuordnung)
s = neu("Slide")
B.kopf(s, "Prüfung", "Derselbe Umsatz auf drei unabhängigen Wegen: 14.522.378,70 Euro", QUELLE)
B.zuordnung(s, Y0, [
    ("Weg 1 · SQL", "SUM(net_total) aus fact_orders — der direkte Weg"),
    ("Weg 2 · OBT", "SUM(net_total) aus obt_orders — andere Datei, andere Struktur"),
    ("Weg 3 · pandas", "SUM(gross_total − discount_amount) — aus den Bestandteilen rekonstruiert"),
], rh=44.0, gap=10.0)
B.band(s, Y0 + 175, 62, [
    "Der dritte Weg ist der wertvollste: Er summiert nicht dieselbe Spalte noch einmal, sondern "
    "baut den Wert aus seinen Bestandteilen neu auf. Er würde auch dann anschlagen, wenn "
    "net_total selbst fehlerhaft befüllt wäre."])

# --- 13 Prüflauf (Tool_Slide)
s = neu("Tool_Slide")
B.kopf(s, "Prüfung", "Ein Skript rechnet 79 dokumentierte Angaben nach", QUELLE)
einl(s, ["Eine Dokumentation, die Zahlen nennt, veraltet still. Wird der Datenbestand geändert, "
         "bleibt der Text stehen und wird unbemerkt falsch. Ein Prüflauf macht die Zusage "
         "wiederholbar: Er rechnet jede Angabe neu und endet mit einem Fehlercode, sobald eine "
         "davon nicht mehr stimmt."], variante="Tool_Slide")
B.kachel(s, CL, Y0, w2, 130, METHOD, "Aufruf",
         ["python verify_readme.py",
          "Exit-Code 0 — alle Angaben bestätigt",
          "Exit-Code 1 — mindestens eine Abweichung, einzeln benannt",
          "Laufzeit rund eine Minute"])
B.kachel(s, CL + w2 + 18, Y0, w2, 130, GOOD, "Was geprüft wird",
         ["Zeilenzahlen exakt, gerundete Kennzahlen mit Toleranz",
          "alle Muster: Saisonalität, Kanäle, Zahlarten, Warenkorb",
          "aktueller Stand: 79 von 79 bestätigt"])
B.band(s, Y0 + 148, 54, [
    "Dieselbe Idee trägt die ETL-Strecke: Die One Big Table wird byte-identisch reproduziert — "
    "ein schärferes Kriterium als \"die Zahlen stimmen\"."])

# --- 14 Prüfbefunde (klammer)
s = neu("Lisa_Slide")
B.kopf(s, "Prüfung", "Die gründlichere Prüfung fand mehr Fehler — auch in sich selbst", QUELLE)
einl(s, ["Der Bericht wurde zweimal vollständig geprüft. Dass der zweite Durchgang mehr Fehler "
         "fand als der erste, ist kein Widerspruch, sondern der Normalfall: Die erste Runde prüft "
         "die auffälligen Werte, die zweite alle. Aufschlussreich ist, womit der zweite Bericht "
         "beginnt — nicht mit Zahlen, sondern mit Wertebereichen."])
B.klammer(s, CL, Y0, w2, 140, "Was der Schema-Abgleich aufdeckte",
          ["Es gibt keinen Kanal \"Online\" — der Bericht verwendete einen Namen, den die Daten nicht kennen",
           "promo_id = 0 bedeutet \"keine Promotion\": wer Nicht-NULL zählt, erhält 100 statt 8,2 Prozent",
           "loyalty_tier \"None\" ist eine Zeichenkette, kein fehlender Wert"], farbe=BLUE)
B.klammer(s, CL + w2 + 18, Y0, w2, 140, "Wo die Prüfung selbst danebenlag",
          ["LIKE '%Cola%' trifft auch Milkshake Chocolate und Hot Chocolate",
           "LIKE '%BBQ%' trifft neben der BBQ Sauce den Burger BBQ Smokehouse",
           "Ein Filter über Textmuster ist eine Annahme über Benennung, keine fachliche Auswahl"], farbe=WARN)

# --- 15 Kette (chevron_kette)
s = neu("Slide")
B.kopf(s, "Auslieferung", "Von der Rohdatei bis zur veröffentlichten Seite", QUELLE)
B.chevron_kette(s, Y0 + 20, 66, [
    "Quellsysteme", "Galaxy-Schema", "One Big Table", "Prüflauf",
    "Pages-Deploy"], farbe=BLUE)
B.kachel(s, CL, Y0 + 122, w2, 182, METHOD, "Was reproduzierbar ist",
         ["generate_obt.py erzeugt die One Big Table byte-identisch — 185.011.332 Bytes, gleiche Prüfsumme",
          "load_duckdb.py lädt den Bestand lokal oder in die Cloud und prüft jede Zeilenzahl",
          "Die Dokumentation liegt im Repository und wird mitversioniert",
          "Der Pages-Ablauf liefert nur das Webverzeichnis aus: rund 530 Kilobyte statt 325 Megabyte"])
B.kachel(s, CL + w2 + 18, Y0 + 122, w2, 182, HINT, "Was offen bleibt",
         ["Kein inkrementelles Laden — die OBT wird immer vollständig neu erzeugt",
          "Keine Historisierung: eine geänderte Loyalty-Stufe überschreibt den alten Wert",
          "Die Interpretationstexte im Bericht sind nicht geprüft, nur die Zahlen"])

# ══════════════════════════════════════════════ speichern
for sl in prs.slides:
    aufraeumen(sl)
prs.save(str(OUT))
print(f"Geschrieben: {OUT.name}")
print(f"Folien: {len(prs.slides.__iter__.__self__._sldIdLst)}")
if B.WARNINGS:
    print("\nWarnungen der Bausteine:")
    for x in B.WARNINGS:
        print("  -", x)
else:
    print("Keine Baustein-Warnungen.")
